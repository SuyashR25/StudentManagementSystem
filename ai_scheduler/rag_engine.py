import os
from typing import List, Optional
from langchain_community.document_loaders import PyMuPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_pinecone import PineconeVectorStore
from langchain_core.documents import Document
from pinecone import Pinecone, ServerlessSpec
from pptx import Presentation
from llm_config import embeddings
from dotenv import load_dotenv

load_dotenv()

# Supported file extensions
SUPPORTED_EXTENSIONS = {'.pdf', '.ppt', '.pptx'}


class PowerPointLoader:
    """Custom loader for PowerPoint files using python-pptx."""
    
    def __init__(self, file_path: str):
        self.file_path = file_path
    
    def load(self) -> List[Document]:
        """Load and parse PowerPoint file into LangChain Documents."""
        documents = []
        
        try:
            prs = Presentation(self.file_path)
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_content = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                    
                    # Handle tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells if cell.text.strip()
                            )
                            if row_text:
                                slide_content.append(row_text)
                
                # Extract notes if available
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_content.append(f"[Speaker Notes: {notes}]")
                
                if slide_content:
                    content = "\n".join(slide_content)
                    doc = Document(
                        page_content=content,
                        metadata={
                            "source": self.file_path,
                            "slide_number": slide_num,
                            "file_type": "pptx"
                        }
                    )
                    documents.append(doc)
            
            return documents
            
        except Exception as e:
            print(f"Error loading PowerPoint file {self.file_path}: {e}")
            return []


class VectorStoreManager:
    """Manages document ingestion, chunking, and retrieval using Pinecone."""
    
    def __init__(self, index_name: str = "ched-scheduler"):
        self.index_name = index_name
        self.vector_store: Optional[PineconeVectorStore] = None
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=800,
            chunk_overlap=150,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        
        # Initialize Pinecone client
        api_key = os.getenv("PINECONE_API_KEY")
        if not api_key:
            raise ValueError("PINECONE_API_KEY environment variable is not set.")
        
        self.pc = Pinecone(api_key=api_key)
        
        # Get or create the index
        self._ensure_index_exists()
    
    def _ensure_index_exists(self):
        """Create the Pinecone index if it doesn't exist."""
        existing_indexes = [idx.name for idx in self.pc.list_indexes()]
        
        if self.index_name not in existing_indexes:
            print(f"Creating Pinecone index: {self.index_name}")
            self.pc.create_index(
                name=self.index_name,
                dimension=384,  # all-MiniLM-L6-v2 produces 384-dim vectors
                metric="cosine",
                spec=ServerlessSpec(
                    cloud="aws",
                    region="us-east-1"
                )
            )
        
        # Connect to the index
        self.vector_store = PineconeVectorStore(
            index=self.pc.Index(self.index_name),
            embedding=embeddings,
            text_key="text"
        )
    
    def _get_loader(self, file_path: str):
        """Get appropriate loader based on file extension."""
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.pdf':
            return PyMuPDFLoader(file_path)
        elif ext in {'.ppt', '.pptx'}:
            return PowerPointLoader(file_path)
        else:
            print(f"Unsupported file type: {ext}")
            return None
    
    def load_document_text(self, file_path: str) -> str:
        """Load raw text from a document directly."""
        loader = self._get_loader(file_path)
        if not loader:
            return ""
        try:
            docs = loader.load()
            return "\n\n".join([d.page_content for d in docs])
        except Exception as e:
            print(f"Direct load error: {e}")
            return ""

    
    def ingest_documents(self, file_paths: List[str], user_id: str = "default") -> bool:
        """Ingest and index PDF and PowerPoint documents with user namespace."""
        user_id = str(user_id)
        if not file_paths:
            return False
        
        documents = []
        for path in file_paths:
            if not os.path.exists(path):
                print(f"File not found: {path}")
                continue
            
            ext = os.path.splitext(path)[1].lower()
            if ext not in SUPPORTED_EXTENSIONS:
                print(f"Unsupported file type: {ext}. Supported: {SUPPORTED_EXTENSIONS}")
                continue
            
            loader = self._get_loader(path)
            if loader:
                try:
                    docs = loader.load()
                    # Add source and user metadata
                    for doc in docs:
                        doc.metadata["source_file"] = os.path.basename(path)
                        doc.metadata["user_id"] = user_id
                        doc.metadata["file_type"] = ext.replace('.', '')
                    documents.extend(docs)
                    print(f"Successfully loaded: {path} ({len(docs)} pages/slides)")
                except Exception as e:
                    print(f"Error loading {path}: {e}")
        
        if not documents:
            return False
        
        chunks = self.text_splitter.split_documents(documents)
        print(f"Created {len(chunks)} chunks from {len(documents)} documents")
        
        # Add to Pinecone with user namespace
        self.vector_store.add_documents(chunks, namespace=user_id)
        return True
    
    def retrieve(self, query: str, user_id: str = "default", k: int = 5) -> List[tuple]:
        """Retrieve relevant chunks with scores from user's namespace."""
        user_id = str(user_id)
        if not self.vector_store:
            return []
        
        results = self.vector_store.similarity_search_with_score(
            query, 
            k=k,
            namespace=user_id
        )
        return results
    
    def get_retriever(self, user_id: str = "default", k: int = 4):
        """Get retriever for chain usage."""
        if not self.vector_store:
            return None
        return self.vector_store.as_retriever(
            search_type='similarity', 
            search_kwargs={"k": k, "namespace": user_id}
        )
    
    def delete_user_data(self, user_id: str) -> bool:
        """Delete all vectors for a specific user."""
        try:
            index = self.pc.Index(self.index_name)
            index.delete(delete_all=True, namespace=user_id)
            return True
        except Exception as e:
            print(f"Error deleting user data: {e}")
            return False
    
    def get_supported_extensions(self) -> set:
        """Return set of supported file extensions."""
        return SUPPORTED_EXTENSIONS


# Global vector store manager
vector_manager = VectorStoreManager()
